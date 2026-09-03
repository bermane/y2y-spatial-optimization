"""director_core.py -- helpers for the Y2Y DIRECTOR PACKAGE (Gate 5 deliverable).

Spec: `analyses/y2y/spec/director_package_spec.md` v1.1 (subordinate to the study plan v0.14.1).
Consumed by `19_director_surfaces.ipynb` (surfaces, clustering, tables, GeoTIFFs) and
`20_director_figures.ipynb` (hex choropleths, cluster overlays, star grid, Act-3 map, E17 one-pager,
deck). Presentation decisions live in the package spec, methods decisions in the study plan; the
pre-stated constants below are copied from the spec verbatim and logged in methods_log.

Semantics: every surface here is the GUARDED band (per-block capture floors capture_b >= 0.95 x
anchor_b appended to the 5% band; spec v0.14 ruling = applied headline). The unguarded band is
carried alongside only for the T-D2 side-by-side (the E15 doubling slide). The estimator is the
hierarchical mean used by 13_gate4_analysis: f_s = mean over {anchor + 50 members}, F = mean over
formulations (one vote each).
"""
from pathlib import Path
from types import SimpleNamespace
import json
import math

import numpy as np
import pandas as pd
import rasterio
from rasterio import features as rfeatures
from scipy import ndimage
from pyproj import Transformer
import geopandas as gpd
from shapely.geometry import Polygon, Point, shape
from shapely.ops import unary_union

import config
import leverage_core as lc
import ensemble_core as ec

ROOT = Path(config.PROJECT_DIR)
Y2Y = ROOT / "analyses" / "y2y"
SPEC = Y2Y / "spec"
RUNS = Y2Y / "runs"
PKG = Y2Y / "director_package"

# ---- pre-stated constants (director_package_spec.md v1.1) ------------------------------------
FREQ_THR = 0.70            # pre-registered frequent band (decision b)
SENS_THRS = (0.60, 0.80)   # sensitivity companion thresholds (appendix)
NEVER_THR = 0.05
MIN_KM2 = 100              # decision (d); deliberately NOT config.CLUSTER_MIN_CELLS (the 04-era 25)
CLOSE_R = 1                # morphological closing radius in cells (bridge single-cell speckle only)
HEX_KM2 = 250              # decision (c) default
HEX_KM2_ALT = 800          # board-level legibility variant, rendered for comparison
POOL_JACCARD_MIN = 0.80    # decision (g): pool the two climate levels unless their frequent tiers diverge
TOPK_ACT1 = 6              # deck shows top-k by area (tie-break mean guarded F); the register ships in full
TOPK_ACT2 = 2
FLOOR_G = 0.05
RARE_EFG_PCT = 0.01       # "rarest" EFG companion mask: presence <= 1% of the PU (disclosed alongside the spec mask)
BANDS = [("never", 0.0, NEVER_THR), ("rare", NEVER_THR, 0.30), ("conditional", 0.30, FREQ_THR),
         ("frequent", FREQ_THR, 0.95), ("always", 0.95, 1.0001)]
# block star axes: member percentiles are averaged with these weights (carbon = mass split 74.2/25.8,
# the Gate-1 measurement); representativeness is a DIFFERENT construction (EFG classes present / 40)
BLOCK_AXES = {
    "core habitat":  {"climate_type_macrorefugia": 1.0},
    "connectivity":  {"transboundary_connectivity": 0.5, "climate_corridors": 0.5},
    "biodiversity":  {"aoh_richness_birds": 0.5, "aoh_richness_mammals": 0.5},
    "carbon":        {"irrecoverable_carbon_m_soc": 0.742, "irrecoverable_carbon_biomass": 0.258},
    "intactness":    {"human_modification": 1.0},     # disclosed, not a driver (dashed/grey)
}
STAR_AXES = ["core habitat", "connectivity", "biodiversity", "carbon", "representativeness", "intactness"]
N_EFG = 40
SCENARIO_LABEL = {"s0": "Balanced", "s1": "Core-habitat-forward", "s2": "Connectivity-forward",
                  "s3": "Biodiversity-forward", "s4": "Carbon-forward", "s5": "Intactness push (S0 + gHM x10)",
                  "s1x": "Core-habitat x carbon regime", "s3x": "Biodiversity x carbon regime"}
ACT2_SCENARIOS = ["s1", "s2", "s3", "s4"]   # the four named forward scenarios (spec Act 2)
SCENARIO_STATEMENT = {
    "s0": "all four value themes hold their intended influence shares",
    "s1": "climate macrorefugia (core habitat) carries a doubled influence share",
    "s2": "connectivity (transboundary current + climate corridors) carries a doubled share",
    "s3": "AOH richness (birds + mammals) carries a doubled share",
    "s4": "carbon carries a doubled share and the mineral-soil target rises 0.332 -> 0.552 (theta 3x)",
    "s5": "S0 with the (inexpressible) intactness layer pushed x10 -- the Claim-B demonstration",
    "s1x": "S1's shares under the carbon-forward target regime (regime flipped alone)",
    "s3x": "S3's shares under the carbon-forward target regime (regime flipped alone)",
}
# decision (h): Nations' own DECLARED proposals only, never analyst-drawn. Indigenous-led rows of the
# corridor-wide proposed-PA file: IPCA-typed + Indigenous-governed (GOV_TYPE 5) + the Ross River NPR
# proposal (Kaska-led; the 04a "manual area"). Confirmed by Ethan before final render.
IPCA_SPEC = dict(vector=config.PROPOSED_PA_VECTOR, name_field="PA_NAME")
# T-D4 (spec v1.3): tier area by ecozone/ecoregion. No such layer is in input_data yet -- drop a vector
# (e.g. CEC North American Level II/III ecoregions, seamless US+Canada) into this folder and 19 picks it
# up; until then the T-D4 cell reports itself PENDING rather than failing the run.
ECOREGIONS_DIR = config.INPUT_DIR / "ecoregions"
ECOREGION_NAME_FIELDS = ("NA_L2NAME", "NA_L3NAME", "ECOZONE_NAME", "ZONE_NAME", "ECOREGION", "REGION_NAM", "NAME", "name")
def ipca_rule(df):
    return (df["PA_TYPE"].astype(str).str.strip() == "IPCA") | (df["GOV_TYPE"].astype(str).str.strip() == "5") \
        | df["PA_NAME"].astype(str).str.contains("Ross River")


def ensure_dirs(pkg=PKG):
    for sub in ("geotiffs", "tables", "figures"):
        (pkg / sub).mkdir(parents=True, exist_ok=True)
    return pkg


# ---- grid ------------------------------------------------------------------------------------
def grid():
    """The 1 km analysis grid + PU / locked / discretionary masks (1-D vectors run over PU cells)."""
    with rasterio.open(config.HANDOFF_DIR / "cost_uniform.tif") as src:
        tr, shp, prof = src.transform, src.shape, src.profile
    pu = lc.pu_mask()
    with rasterio.open(config.HANDOFF_DIR / "mask_protected_areas.tif") as src:
        locked2d = (src.read(1) == 1) & pu
    locked = locked2d[pu]
    G = SimpleNamespace(pu=pu, locked2d=locked2d, locked=locked, disc=~locked, n_pu=int(pu.sum()),
                        n_disc=int((~locked).sum()), shape=shp, transform=tr, crs=config.TARGET_CRS,
                        profile=prof, cell_km2=abs(tr.a * tr.e) / 1e6)
    G.rows, G.cols = np.where(pu)
    return G


def to_grid(G, v, fill=np.nan, dtype=np.float32):
    out = np.full(G.shape, fill, dtype=dtype)
    out[G.pu] = v
    return out


def latlon(G):
    """Per-PU-cell latitude/longitude (cached on G)."""
    if not hasattr(G, "lat"):
        xs = G.transform.c + (G.cols + 0.5) * G.transform.a
        ys = G.transform.f + (G.rows + 0.5) * G.transform.e
        lon, lat = Transformer.from_crs(G.crs, "EPSG:4326", always_xy=True).transform(xs, ys)
        G.lat, G.lon = lat.astype(np.float32), lon.astype(np.float32)
    return G.lat, G.lon


def write_tif(G, v, path, dtype="float32", nodata=np.nan):
    arr = v if v.ndim == 2 else to_grid(G, v, fill=nodata, dtype=dtype)
    prof = G.profile | dict(dtype=dtype, count=1, nodata=nodata, compress="deflate", tiled=True)
    with rasterio.open(path, "w", **prof) as dst:
        dst.write(arr.astype(dtype), 1)
    return path


# ---- loading the sweeps ----------------------------------------------------------------------
def _diam(S, disc, m_disc):
    Sd = S[:, disc].astype(np.float32)
    sizes = Sd.sum(axis=1)
    ham = sizes[:, None] + sizes[None, :] - 2 * (Sd @ Sd.T)
    return float(ham.max() / (2 * m_disc))


def load_guarded(G, MAN, allow_partial=False, diameters=True):
    """Stream every formulation once: guarded + plain f, anchors, band unions, diameters, certificates.

    allow_partial=True is a DEV flag (smoke-runs on the S0/S4 artifacts from notebook 16); the
    production run asserts all 14 guarded sweeps exist (director spec step 0)."""
    L = SimpleNamespace(forms=[], f_guard={}, f_plain={}, anchors={}, union_guard={}, union_plain={},
                        D_guard={}, D_plain={}, cert={}, meta={}, f_maa_guard={})
    missing = []
    for _, row in MAN.iterrows():
        fid = row.formulation_id
        cd = RUNS / fid
        need = [cd / "anchor.tif", cd / "mga_g05.tif", cd / "mga_guard_g05.tif", cd / "certificates_guard.csv"]
        if not all(p.exists() for p in need):
            missing.append(fid)
            continue
        cert = pd.read_csv(cd / "certificates_guard.csv")
        assert bool(cert.band_ok.all()), f"{fid}: guarded band certificate violated"
        assert len(cert) == int(row.k_requested), f"{fid}: {len(cert)} guarded members, expected {row.k_requested}"
        A = ec.read_selections(cd / "anchor.tif", G.pu)[0]
        Sg = np.vstack([A[None, :], ec.read_selections(cd / "mga_guard_g05.tif", G.pu)])
        Sp = np.vstack([A[None, :], ec.read_selections(cd / "mga_g05.tif", G.pu)])
        m_disc = int(A[G.disc].sum())
        L.f_guard[fid] = Sg.mean(axis=0).astype(np.float32)
        L.f_plain[fid] = Sp.mean(axis=0).astype(np.float32)
        L.union_guard[fid] = Sg.any(axis=0)
        L.union_plain[fid] = Sp.any(axis=0)
        L.anchors[fid] = A
        if diameters:
            L.D_guard[fid] = _diam(Sg, G.disc, m_disc)
            L.D_plain[fid] = _diam(Sp, G.disc, m_disc)
        L.cert[fid] = dict(n=len(cert), dup=int(cert.duplicate.sum()), runtime_min=float(cert.runtime_s.sum() / 60),
                           time_limited=int((cert.status == "TIME_LIMIT").sum()))
        L.meta[fid] = json.loads((cd / "formulation_meta.json").read_text())
        maa = cd / "maa_guard_g05.tif"
        if maa.exists():
            Sm = np.vstack([A[None, :], ec.read_selections(maa, G.pu)])
            L.f_maa_guard[fid] = Sm.mean(axis=0).astype(np.float32)
            del Sm
        L.forms.append(fid)
        del Sg, Sp
        print(f"{fid:<22} f_guard freq {int((L.f_guard[fid][G.disc] >= FREQ_THR).sum()):>7,} km2 | "
              f"plain {int((L.f_plain[fid][G.disc] >= FREQ_THR).sum()):>7,} km2"
              + (f" | D {L.D_plain[fid]:.3f} -> {L.D_guard[fid]:.3f}" if diameters else ""))
    if missing:
        msg = f"guarded sweep missing for {len(missing)} formulation(s): {missing} -- run 18_guarded_sweep first"
        if not allow_partial:
            raise FileNotFoundError(msg)
        print("ALLOW_PARTIAL:", msg)
    L.missing = missing
    return L


def ensemble(fdict, forms):
    """Hierarchical F: one vote per formulation (13_gate4_analysis's estimator)."""
    return np.mean([fdict[c] for c in forms], axis=0).astype(np.float32)


def union_membership(L, forms, guarded=True):
    U = L.union_guard if guarded else L.union_plain
    return np.mean([U[c] for c in forms], axis=0).astype(np.float32)


# ---- tiers ------------------------------------------------------------------------------------
def band_masks(F):
    return {name: (F >= lo) & (F < hi) for name, lo, hi in BANDS}


def band_table(G, surfaces):
    """T-D2 part 1: spec bands x semantics, discretionary km2 and % (surfaces = {label: F})."""
    rows = []
    for name, lo, hi in BANDS:
        r = {"band": f"{name} [{lo:.2f}, {min(hi, 1.0):.2f}{')' if hi < 1 else ']'}"}
        for lab, F in surfaces.items():
            km2 = int(((F >= lo) & (F < hi) & G.disc).sum())
            r[f"{lab} km2"] = km2
            r[f"{lab} %disc"] = 100 * km2 / G.n_disc
        rows.append(r)
    return pd.DataFrame(rows)


def jaccard(a, b):
    u = int((a | b).sum())
    return float((a & b).sum() / u) if u else float("nan")


def pool_scenarios(G, fdict, MAN, thr=FREQ_THR, jmin=POOL_JACCARD_MIN):
    """Decision (g): per scenario, pool the two climate levels unless their frequent tiers diverge.

    Returns (POOL {key: f}, report df). key = scenario_id when pooled, '<sid>@<level>' otherwise."""
    POOL, rep = {}, []
    for sid, grp in MAN.groupby("scenario_id", sort=False):
        fids = [f for f in grp.formulation_id if f in fdict]
        if not fids:
            continue
        if len(fids) == 1:
            POOL[sid] = fdict[fids[0]]
            rep.append(dict(scenario=sid, levels=1, jaccard=np.nan, decision="single level",
                            **{f"freq_km2_{i}": int((fdict[f][G.disc] >= thr).sum()) for i, f in enumerate(fids)}))
            continue
        tiers = [(fdict[f] >= thr) & G.disc for f in fids]
        J = jaccard(tiers[0], tiers[1])
        if J >= jmin:
            POOL[sid] = np.mean([fdict[f] for f in fids], axis=0).astype(np.float32)
            dec = "POOLED"
        else:
            for f in fids:
                POOL[f"{sid}@{'245' if 'ssp245' in f else '585'}"] = fdict[f]
            dec = "SEPARATE (diverged)"
        rep.append(dict(scenario=sid, levels=len(fids), jaccard=J, decision=dec,
                        freq_km2_585=int(tiers[0].sum()), freq_km2_245=int(tiers[1].sum())))
    return POOL, pd.DataFrame(rep)


# ---- clustering (pre-stated procedure, spec v1.1) --------------------------------------------
def clusters(G, surf, thr=FREQ_THR, min_km2=MIN_KM2, close_r=CLOSE_R, subtract2d=None):
    """Threshold -> closing(r) -> 8-connected components -> register (all components; kept = >= min).

    subtract2d: Act-1 core footprint to remove AFTER clustering (Act 2), with % overlap reported."""
    F2 = to_grid(G, surf, fill=0.0)
    m = (F2 >= thr) & G.pu & ~G.locked2d
    if close_r:
        m = ndimage.binary_closing(m, structure=np.ones((3, 3), bool), iterations=close_r) & G.pu & ~G.locked2d
        m |= (F2 >= thr) & G.pu & ~G.locked2d          # closing never removes original cells
    lab, n = ndimage.label(m, structure=np.ones((3, 3), int))
    if n == 0:
        return lab, pd.DataFrame(columns=["cid", "cells", "km2", "meanF", "minF", "kept"])
    ids = np.arange(1, n + 1)
    cells = ndimage.sum(m, lab, ids).astype(int)
    meanF = ndimage.mean(F2, lab, ids)
    minF = ndimage.minimum(np.where(m, F2, 9.0), lab, ids)
    cy, cx = zip(*ndimage.center_of_mass(m, lab, ids))
    reg = pd.DataFrame(dict(cid=ids, cells=cells, km2=cells * G.cell_km2, meanF=meanF, minF=minF,
                            row=np.array(cy), col=np.array(cx)))
    if subtract2d is not None:
        ov = ndimage.sum(m & subtract2d, lab, ids).astype(int)
        reg["core_overlap_pct"] = 100 * ov / cells
        reg["residual_km2"] = (cells - ov) * G.cell_km2
        reg["kept"] = reg.residual_km2 >= min_km2
    else:
        reg["kept"] = reg.km2 >= min_km2
    x = G.transform.c + (reg.col + 0.5) * G.transform.a
    y = G.transform.f + (reg.row + 0.5) * G.transform.e
    lon, lat = Transformer.from_crs(G.crs, "EPSG:4326", always_xy=True).transform(x.values, y.values)
    reg["lat"], reg["lon"] = lat, lon
    reg = reg.sort_values(["km2", "meanF"], ascending=False).reset_index(drop=True)
    return lab, reg


def top_k(reg, k):
    return reg[reg.kept].sort_values(["km2", "meanF"], ascending=False).head(k)


def sensitivity(G, surf, thrs=(SENS_THRS[0], FREQ_THR, SENS_THRS[1]), **kw):
    rows = []
    for t in thrs:
        _, reg = clusters(G, surf, thr=t, **kw)
        kept = reg[reg.kept] if len(reg) else reg
        rows.append(dict(threshold=t, tier_km2=int(((surf >= t) & G.disc).sum()),
                         n_components=len(reg), n_kept=len(kept), kept_km2=float(kept.km2.sum()) if len(kept) else 0.0,
                         largest_km2=float(kept.km2.max()) if len(kept) else 0.0))
    return pd.DataFrame(rows)


def mask_of(lab, cid):
    return lab == cid


def vectorize(G, lab, ids, simplify_m=2000):
    """Cluster polygons (topology-preserving simplification, ~2 km tolerance) as a GeoDataFrame."""
    recs = []
    arr = np.where(np.isin(lab, ids), lab, 0).astype(np.int32)
    geoms = {}
    for geom, val in rfeatures.shapes(arr, mask=arr > 0, transform=G.transform):
        geoms.setdefault(int(val), []).append(shape(geom))
    for cid in ids:
        g = unary_union(geoms.get(int(cid), []))
        recs.append(dict(cid=int(cid), geometry=g.simplify(simplify_m, preserve_topology=True)))
    return gpd.GeoDataFrame(recs, geometry="geometry", crs=G.crs)


# ---- hex aggregation (presentation only; clustering ALWAYS runs at 1 km) ---------------------
def hex_grid(G, area_km2=HEX_KM2):
    """Flat-top hexagon lattice in the analysis CRS covering the PU extent; rasterized to hex ids."""
    R = math.sqrt(2 * area_km2 * 1e6 / (3 * math.sqrt(3)))      # circumradius (m)
    w, h = 2 * R, math.sqrt(3) * R
    x0 = G.transform.c; y1 = G.transform.f
    x1 = x0 + G.shape[1] * G.transform.a; y0 = y1 + G.shape[0] * G.transform.e
    polys, ids = [], []
    nx = int((x1 - x0) / (1.5 * R)) + 3
    ny = int((y1 - y0) / h) + 3
    k = 0
    for i in range(-1, nx):
        cx = x0 + i * 1.5 * R
        for j in range(-1, ny):
            cy = y0 + j * h + (h / 2 if i % 2 else 0)
            k += 1
            ids.append(k)
            polys.append(Polygon([(cx + R * math.cos(a), cy + R * math.sin(a))
                                  for a in np.arange(0, 2 * math.pi, math.pi / 3)]))
    lab = rfeatures.rasterize(zip(polys, ids), out_shape=G.shape, transform=G.transform, fill=0, dtype="int32")
    lab[~G.pu] = 0
    present = np.unique(lab[lab > 0])
    gdf = gpd.GeoDataFrame({"hex_id": ids, "geometry": polys}, crs=G.crs)
    gdf = gdf[gdf.hex_id.isin(present)].reset_index(drop=True)
    return gdf, lab


def hex_means(G, gdf, lab, surf, disc_only=True):
    """Mean of a PU surface per hex over DISCRETIONARY cells (PAs are drawn as their own layer)."""
    F2 = to_grid(G, surf, fill=np.nan)
    valid = G.pu & (~G.locked2d if disc_only else True) & np.isfinite(F2)
    labv = np.where(valid, lab, 0)
    ids = gdf.hex_id.values
    with np.errstate(invalid="ignore", divide="ignore"):      # hexes with no valid cell -> NaN, by design
        means = ndimage.mean(np.nan_to_num(F2), labv, ids)
    counts = ndimage.sum(valid, labv, ids)
    out = gdf.copy()
    out["value"] = np.where(counts > 0, means, np.nan)
    out["n_disc"] = counts.astype(int)
    out["n_pu"] = ndimage.sum(G.pu, lab, ids).astype(int)
    return out


# ---- cartography helpers (pixel-space maps: 1 px = 1 km) --------------------------------------
def xy_to_px(G, x, y):
    return (x - G.transform.c) / G.transform.a, (y - G.transform.f) / G.transform.e


def graticule(ax, G, lats=(45, 50, 53, 55, 60, 65), lons=(-130, -125, -120, -115, -110), emph_lat=53,
              color="#555555", lw=0.4):
    """Lat/lon graticule in pixel coordinates; 53 N emphasized (ties to the E17 one-pager)."""
    T = Transformer.from_crs("EPSG:4326", G.crs, always_xy=True)
    H, W = G.shape
    for la in lats:
        lo = np.linspace(-145, -95, 200)
        x, y = T.transform(lo, np.full_like(lo, la, dtype=float))
        px, py = xy_to_px(G, x, y)
        e = la == emph_lat
        ax.plot(px, py, color="#b00020" if e else color, lw=1.1 if e else lw, ls="-" if e else ":", zorder=3,
                clip_on=True)
        inside = (px >= 0) & (px < W) & (py >= 0) & (py < H)
        if inside.any():
            i = np.argmax(inside)
            ax.text(px[i] + 4, py[i] - 3, f"{la}°N", fontsize=7 if not e else 8, color="#b00020" if e else color,
                    fontweight="bold" if e else "normal", zorder=4)
    for lo in lons:
        la = np.linspace(38, 72, 200)
        x, y = T.transform(np.full_like(la, lo, dtype=float), la)
        px, py = xy_to_px(G, x, y)
        ax.plot(px, py, color=color, lw=lw, ls=":", zorder=3, clip_on=True)
    ax.set_xlim(0, W); ax.set_ylim(H, 0)


def scalebar(ax, G, km=250, loc=(0.74, 0.10)):
    H, W = G.shape
    px_per_km = 1000 / abs(G.transform.a)
    x0, y0 = loc[0] * W, (1 - loc[1]) * H
    ax.plot([x0, x0 + km * px_per_km], [y0, y0], color="black", lw=2.5, solid_capstyle="butt", zorder=5)
    ax.text(x0 + km * px_per_km / 2, y0 - 12, f"{km} km", ha="center", fontsize=8, zorder=5)


def corner_note(ax, text, loc="lower right"):
    ha = "right" if "right" in loc else "left"
    ax.text(0.99 if ha == "right" else 0.01, 0.005, text, transform=ax.transAxes, ha=ha, va="bottom",
            fontsize=7.5, color="#333333", zorder=6)


# ---- block percentiles (star plots) -----------------------------------------------------------
def block_percentiles(G):
    """Per-cell percentile of each hand-off layer over the DISCRETIONARY landscape (0.5 = typical
    unprotected land), combined into the five block axes; EFG presence stack for representativeness."""
    axes = {}
    pct = {}
    feats = sorted({f for d in BLOCK_AXES.values() for f in d})
    for f in feats:
        v = np.nan_to_num(lc._read(config.HANDOFF_DIR / f"{f}.tif")[G.pu], nan=0.0)
        ref = np.sort(v[G.disc])
        pct[f] = (np.searchsorted(ref, v, side="right") / len(ref)).astype(np.float32)
    for ax, members in BLOCK_AXES.items():
        axes[ax] = sum(w * pct[f] for f, w in members.items()).astype(np.float32)
    paths = lc.efg_paths()
    efg = np.zeros((len(paths), G.n_pu), bool)
    for i, p in enumerate(paths):
        efg[i] = np.nan_to_num(lc._read(p)[G.pu], nan=0.0) > 0
    return SimpleNamespace(axes=axes, pct=pct, efg=efg, efg_names=[p.stem for p in paths])


def star_profile(P, mask1d):
    out = {ax: float(P.axes[ax][mask1d].mean()) for ax in BLOCK_AXES}
    out["representativeness"] = float(P.efg[:, mask1d].any(axis=1).sum() / P.efg.shape[0])
    return {ax: out[ax] for ax in STAR_AXES}


# ---- driver attribution masks (E13 definitions, notebook 15) ----------------------------------
def driver_masks(G):
    theta = config.AUDIT["theta"]
    rare_cap = config.AUDIT["rare_cap"]
    v = lc._read(config.HANDOFF_DIR / "irrecoverable_carbon_m_soc.tif")
    masks = {"m_soc theta-tail": np.nan_to_num(v, nan=-1)[G.pu] >= theta * float(np.nanmean(v[G.pu]))}
    conn = lc._read(config.HANDOFF_DIR / "transboundary_connectivity.tif")[G.pu]
    masks["connectivity spike (top 0.2%)"] = conn >= np.nanquantile(conn, 0.998)
    rare = np.zeros(G.n_pu, bool)
    n_rare = 0
    for p in lc.efg_paths():
        e = np.nan_to_num(lc._read(p)[G.pu], nan=0.0)
        cap_max = lc.leverage_of(e)[1]
        if cap_max >= rare_cap:            # rare-attainable = fits inside the budget entirely (Gate 0a rule)
            rare |= e > 0
            n_rare += 1
    masks["rare-attainable EFG footprint"] = rare
    # The spec's "rare-EFG footprint" is ambiguous: rare-ATTAINABLE (fits in the budget; 36/40) unions to
    # ~79% of the region and cannot attribute anything, so a discriminating companion is reported too:
    # EFGs whose presence covers <= RARE_EFG_PCT of the PU (the genuinely scarce classes). Both disclosed.
    rarest = np.zeros(G.n_pu, bool)
    n_rarest = 0
    for p in lc.efg_paths():
        e = np.nan_to_num(lc._read(p)[G.pu], nan=0.0) > 0
        if e.sum() <= RARE_EFG_PCT * G.n_pu:
            rarest |= e
            n_rarest += 1
    masks[f"rarest-EFG footprint (<= {100 * RARE_EFG_PCT:g}% of PU each)"] = rarest
    print(f"driver masks: m_soc theta-tail {int(masks['m_soc theta-tail'].sum()):,} cells | spike "
          f"{int(masks['connectivity spike (top 0.2%)'].sum()):,} | rare-attainable EFG footprint "
          f"{int(rare.sum()):,} cells ({100 * rare.mean():.0f}% of PU) from {n_rare}/{len(lc.efg_paths())} EFGs | "
          f"rarest-EFG footprint {int(rarest.sum()):,} cells ({100 * rarest.mean():.1f}%) from {n_rarest} EFGs")
    return masks


# ---- vector context: proposed IPCAs, existing PAs, placeholder names --------------------------
def ipca_layer(G):
    g = gpd.read_file(IPCA_SPEC["vector"]).to_crs(G.crs)
    g = g[ipca_rule(g)].reset_index(drop=True)
    g["name"] = g[IPCA_SPEC["name_field"]].astype(str)
    mask2d = rfeatures.rasterize(((geom, 1) for geom in g.geometry), out_shape=G.shape,
                                 transform=G.transform, fill=0, dtype="uint8").astype(bool) & G.pu
    return SimpleNamespace(gdf=g, mask2d=mask2d)


def pa_layer(G, min_km2=300):
    pa = gpd.read_file(config.PA_VECTOR).to_crs(G.crs).dissolve(by="PA_Name").reset_index()
    pa["km2"] = pa.geometry.area / 1e6
    return pa[pa.km2 >= min_km2].reset_index(drop=True)


def _bearing(dx, dy):
    ang = (math.degrees(math.atan2(dx, dy)) + 360) % 360
    return ["N", "NE", "E", "SE", "S", "SW", "W", "NW"][int((ang + 22.5) // 45) % 8]


def placeholder_name(G, comp_mask2d, named):
    """'near <area> (<bearing>)' from the nearest named PA/IPCA proposal -- Ethan renames (decision e)."""
    rows, cols = np.where(comp_mask2d)
    cx = G.transform.c + (cols.mean() + 0.5) * G.transform.a
    cy = G.transform.f + (rows.mean() + 0.5) * G.transform.e
    pt = Point(cx, cy)
    d = named.geometry.distance(pt)
    i = int(d.idxmin())
    nm = named.loc[i, "name"]
    lon, lat = Transformer.from_crs(G.crs, "EPSG:4326", always_xy=True).transform(cx, cy)
    tag = f" ({lat:.1f}°N {abs(lon):.1f}°W)"          # keeps names unique when two clusters share a neighbour
    if d[i] == 0:
        return f"{nm} vicinity{tag}"
    c = named.loc[i, "geometry"].centroid
    return f"{_bearing(cx - c.x, cy - c.y)} of {nm}{tag}"


def named_areas(G):
    pa = pa_layer(G)[["PA_Name", "geometry"]].rename(columns={"PA_Name": "name"})
    ip = ipca_layer(G).gdf[["name", "geometry"]]
    return gpd.GeoDataFrame(pd.concat([pa, ip], ignore_index=True), geometry="geometry", crs=G.crs)


# ---- E17 one-pager inputs ----------------------------------------------------------------------
def e17_shifts(G):
    """Leave-one-block-out latitude shifts vs the S0 anchor (from runs/e17_t3, notebook 16)."""
    lat, _ = latlon(G)
    s0 = ec.read_selections(RUNS / "s0_ssp585_theta5" / "anchor.tif", G.pu)[0]
    base = float(lat[s0 & G.disc].mean())
    rows = []
    for b in ["core_habitat", "connectivity", "biodiversity", "carbon", "efg"]:
        p = RUNS / "e17_t3" / f"{b}_out" / "run" / "portfolio.tif"
        if not p.exists():
            continue
        sel = ec.read_selections(p, G.pu)[0]
        rows.append(dict(block_out=b, mean_lat=float(lat[sel & G.disc].mean()),
                         delta_lat=float(lat[sel & G.disc].mean() - base), jaccard_vs_s0=jaccard(sel, s0)))
    return base, pd.DataFrame(rows)


# ---- star grid + deck --------------------------------------------------------------------------
def plot_star_grid(profiles, path, title, ncols=4, rmax=1.0, ref=0.5):
    """profiles: list of dict(title=, values={axis: v}, color=). One shared radial scale."""
    import matplotlib.pyplot as plt
    n = len(profiles)
    ncols = min(ncols, max(n, 1))
    nrows = int(math.ceil(n / ncols))
    fig, axes = plt.subplots(nrows, ncols, figsize=(3.6 * ncols, 4.0 * nrows), subplot_kw=dict(polar=True))
    fig.subplots_adjust(wspace=0.55, hspace=0.6)
    axes = np.atleast_1d(axes).ravel()
    k = len(STAR_AXES)
    ang = np.linspace(0, 2 * np.pi, k, endpoint=False)
    for ax, pr in zip(axes, profiles):
        vals = [pr["values"][a] for a in STAR_AXES]
        closed = np.r_[vals, vals[0]]
        ax.plot(np.r_[ang, ang[0]], closed, color=pr.get("color", "#2b4f7d"), lw=1.6)
        ax.fill(np.r_[ang, ang[0]], closed, color=pr.get("color", "#2b4f7d"), alpha=0.25)
        ax.plot(np.r_[ang, ang[0]], [ref] * (k + 1), color="#888888", lw=0.8, ls="--")
        i_int = STAR_AXES.index("intactness")
        ax.plot([ang[i_int], ang[i_int]], [0, rmax], color="#999999", lw=1.0, ls=(0, (2, 2)))
        ax.set_xticks(ang)
        ax.set_xticklabels([a if a != "intactness" else "intactness†" for a in STAR_AXES], fontsize=7.5)
        for lab_, a in zip(ax.get_xticklabels(), STAR_AXES):
            if a == "intactness":
                lab_.set_color("#888888")
        ax.set_ylim(0, rmax)
        ax.set_yticks([0.25, 0.5, 0.75, 1.0]); ax.set_yticklabels(["", "0.5", "", "1"], fontsize=6.5)
        import textwrap
        ax.set_title("\n".join(textwrap.fill(t, 30) for t in pr["title"].split("\n")), fontsize=8, pad=10)
    for ax in axes[n:]:
        ax.axis("off")
    fig.suptitle(title, fontsize=12, y=1.0)
    fig.text(0.01, -0.01, "axes = cluster mean percentile vs the DISCRETIONARY landscape (dashed ring 0.5 = typical "
             "unprotected land); representativeness = EFG classes present / 40 (different construction); "
             "† intactness disclosed, not a driver", fontsize=7, color="#444444")
    fig.savefig(path, dpi=200, bbox_inches="tight")
    return fig


def build_deck(slides, path, subtitle=""):
    """Draft .pptx: one slide per dict(title=, image=, bullets=[], notes=). Editable starting point."""
    from pptx import Presentation
    from pptx.util import Inches, Pt
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
    for s in slides:
        sl = prs.slides.add_slide(prs.slide_layouts[5])
        sl.shapes.title.text = s["title"]
        sl.shapes.title.text_frame.paragraphs[0].font.size = Pt(26)
        img = s.get("image")
        if img and Path(img).exists():
            pic = sl.shapes.add_picture(str(img), Inches(0.4), Inches(1.3), height=Inches(5.9))
            if pic.width > Inches(8.6):
                pic.width, pic.height = Inches(8.6), int(pic.height * Inches(8.6) / pic.width)
        tb = sl.shapes.add_textbox(Inches(9.2), Inches(1.3), Inches(3.9), Inches(5.9)).text_frame
        tb.word_wrap = True
        for i, b in enumerate(s.get("bullets", [])):
            p = tb.paragraphs[0] if i == 0 else tb.add_paragraph()
            p.text = b; p.font.size = Pt(13)
        if s.get("notes"):
            sl.notes_slide.notes_text_frame.text = s["notes"]
    prs.save(str(path))
    return path


# ---- T-D4: ecoregion layer (optional data dependency) -----------------------------------------
def ecoregion_layer(G):
    """First vector in input_data/ecoregions/ rasterized to 1 km zone ids; None when absent (T-D4 pending)."""
    if not ECOREGIONS_DIR.exists():
        return None
    files = sorted([*ECOREGIONS_DIR.glob("*.shp"), *ECOREGIONS_DIR.glob("*.gpkg"), *ECOREGIONS_DIR.glob("*.geojson")])
    if not files:
        return None
    g = gpd.read_file(files[0]).to_crs(G.crs)
    fld = next((f for f in ECOREGION_NAME_FIELDS if f in g.columns), None)
    if fld is None:
        raise ValueError(f"{files[0].name}: no recognised name field among {ECOREGION_NAME_FIELDS}")
    g = g.dissolve(by=fld).reset_index()
    g["zone_id"] = np.arange(1, len(g) + 1)
    zones = rfeatures.rasterize(zip(g.geometry, g.zone_id), out_shape=G.shape, transform=G.transform,
                                fill=0, dtype="int32")
    return SimpleNamespace(gdf=g, name_field=fld, zones=zones, source=files[0].name)


def tier_achievement(G, cumulative_masks):
    """Capture per block for CUMULATIVE tier masks (1-D over PU, locked included) -- zero-solve.

    capture_f = share of feature f's regional total inside the mask; block = mean over its features
    (the T-D3 convention). Returns a long DataFrame (tier, block, feature, capture)."""
    rows = []
    vals = {}
    for b, feats in config.BLOCKS.items():
        for f in feats:
            vals[f] = np.nan_to_num(lc._read(config.HANDOFF_DIR / f"{f}.tif")[G.pu], nan=0.0)
    for tier, m in cumulative_masks.items():
        for b, feats in config.BLOCKS.items():
            caps = [float(vals[f][m].sum() / vals[f].sum()) for f in feats]
            for f, c in zip(feats, caps):
                rows.append(dict(tier=tier, block=b, feature=f, capture=c))
            rows.append(dict(tier=tier, block=b, feature="BLOCK", capture=float(np.mean(caps))))
    return pd.DataFrame(rows)
